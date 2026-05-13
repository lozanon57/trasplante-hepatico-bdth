#!/usr/bin/env python3
"""Genera la presentación BDTH para niños de 5 años."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Paleta ────────────────────────────────────────────────────────────────────
PURPLE   = RGBColor(0x4A, 0x14, 0x8C)
PURPLE_L = RGBColor(0xCE, 0x93, 0xD8)
TEAL     = RGBColor(0x00, 0x6E, 0x6E)
ORANGE   = RGBColor(0xE6, 0x51, 0x00)
RED      = RGBColor(0xC6, 0x28, 0x28)
GREEN    = RGBColor(0x2E, 0x7D, 0x32)
BLUE     = RGBColor(0x15, 0x65, 0xC0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x21, 0x21, 0x21)
LIGHT_BG = RGBColor(0xF3, 0xE5, 0xF5)

W = Inches(13.333)   # widescreen 16:9
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(blank)


def fill_slide(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=40, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, color: RGBColor, radius=False):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Diapositivas ──────────────────────────────────────────────────────────────

def slide_portada(prs):
    """Slide 1 — Portada"""
    s = blank_slide(prs)
    fill_slide(s, PURPLE)

    # Círculo decorativo fondo
    add_rect(s,
             Inches(9.5), Inches(-1), Inches(5), Inches(5),
             RGBColor(0x6A, 0x1B, 0x9A))

    # Emoji grande
    add_text_box(s, "🫀", Inches(0.6), Inches(0.8), Inches(3), Inches(3),
                 font_size=120, align=PP_ALIGN.CENTER)

    # Título
    add_text_box(s, "La libreta mágica\nde los médicos",
                 Inches(3.8), Inches(1.2), Inches(8.8), Inches(2.5),
                 font_size=54, bold=True, color=WHITE)

    # Subtítulo
    add_text_box(s, "Una app que ayuda a los médicos a salvar hígados",
                 Inches(3.8), Inches(3.9), Inches(8.8), Inches(0.8),
                 font_size=22, color=PURPLE_L)

    # Hospital
    add_text_box(s, "Hospital Gregorio Marañón · Madrid",
                 Inches(0.5), Inches(6.7), Inches(8), Inches(0.6),
                 font_size=16, color=PURPLE_L)


def slide_que_es_higado(prs):
    """Slide 2 — ¿Qué es el hígado?"""
    s = blank_slide(prs)
    fill_slide(s, LIGHT_BG)

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), PURPLE)
    add_text_box(s, "¿Qué es el hígado? 🤔",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Emoji hígado
    add_text_box(s, "🫁", Inches(0.3), Inches(1.6), Inches(2.5), Inches(2.5),
                 font_size=110, align=PP_ALIGN.CENTER)

    # Texto explicativo
    items = [
        "🏭  El hígado es como la FÁBRICA de limpieza de tu cuerpo",
        "🧹  Limpia la sangre de cosas malas que comemos o bebemos",
        "💊  Ayuda a que los medicamentos funcionen bien",
        "😢  A veces el hígado se rompe y hay que cambiarlo por uno nuevo",
        "🎁  El nuevo hígado viene de una persona muy generosa: el donante",
    ]
    y = 1.7
    for item in items:
        add_text_box(s, item, Inches(3.2), Inches(y), Inches(9.5), Inches(0.6),
                     font_size=20, color=DARK)
        y += 0.95


def slide_que_es_app(prs):
    """Slide 3 — ¿Qué es la app?"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xE8, 0xF5, 0xE9))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), GREEN)
    add_text_box(s, "¿Qué es la app BDTH? 📱",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s, "📖",
                 Inches(0.5), Inches(1.5), Inches(2), Inches(2),
                 font_size=100, align=PP_ALIGN.CENTER)

    add_text_box(s,
        "Imagina una LIBRETA muy especial\nen el móvil del médico...",
        Inches(3), Inches(1.6), Inches(9.5), Inches(1.0),
        font_size=28, bold=True, color=GREEN)

    items = [
        "✏️   Apunta todo lo que pasa durante el trasplante",
        "🔍   Busca información de cualquier paciente al instante",
        "📊   Hace gráficas solas para aprender qué funciona mejor",
        "🔒   Guarda los datos de forma secreta para proteger la intimidad",
        "✈️   Funciona sin internet, en cualquier sitio del mundo",
    ]
    y = 3.0
    for item in items:
        add_text_box(s, item, Inches(3), Inches(y), Inches(9.8), Inches(0.6),
                     font_size=20, color=DARK)
        y += 0.9


def slide_donante_receptor(prs):
    """Slide 4 — Donante y receptor"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xE3, 0xF2, 0xFD))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), BLUE)
    add_text_box(s, "¿Quién da y quién recibe? 💝",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Donante
    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.6), Inches(5.5),
             RGBColor(0xBB, 0xDE, 0xFB))
    add_text_box(s, "🦸 EL DONANTE",
                 Inches(0.6), Inches(1.6), Inches(5.2), Inches(0.7),
                 font_size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    donante_items = [
        "Una persona muy buena",
        "que ya no necesita su hígado",
        "y lo regala para que otro viva",
        "",
        "La app apunta:",
        "• Su edad y tipo de sangre",
        "• Por qué falleció",
        "• Cómo estaba su hígado",
        "• Cuánto tiempo viajó el hígado",
    ]
    y = 2.5
    for item in donante_items:
        add_text_box(s, item, Inches(0.7), Inches(y), Inches(5.0), Inches(0.45),
                     font_size=17, color=DARK)
        y += 0.42

    # Flecha
    add_text_box(s, "→", Inches(6.1), Inches(3.8), Inches(1.0), Inches(0.8),
                 font_size=60, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s, "🫀", Inches(6.1), Inches(4.4), Inches(1.0), Inches(0.8),
                 font_size=40, align=PP_ALIGN.CENTER)

    # Receptor
    add_rect(s, Inches(7.3), Inches(1.5), Inches(5.6), Inches(5.5),
             RGBColor(0xC8, 0xE6, 0xC9))
    add_text_box(s, "🤒 EL RECEPTOR",
                 Inches(7.5), Inches(1.6), Inches(5.2), Inches(0.7),
                 font_size=24, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    receptor_items = [
        "Una persona enferma",
        "cuyo hígado ya no funciona",
        "y que necesita uno nuevo para vivir",
        "",
        "La app apunta:",
        "• Su enfermedad del hígado",
        "• Qué tan grave está (MELD)",
        "• Cómo fue la operación",
        "• Cómo se recuperó después",
    ]
    y = 2.5
    for item in receptor_items:
        add_text_box(s, item, Inches(7.5), Inches(y), Inches(5.1), Inches(0.45),
                     font_size=17, color=DARK)
        y += 0.42


def slide_como_funciona(prs):
    """Slide 5 — Cómo funciona (pasos)"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xFF, 0xF3, 0xE0))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), ORANGE)
    add_text_box(s, "¿Cómo lo usa el médico? 👨‍⚕️",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    pasos = [
        ("1️⃣", "Llega un donante de hígado",
         "El médico abre la app y crea un caso nuevo con un código secreto"),
        ("2️⃣", "Apunta los datos del donante",
         "Puede escribirlos o hacer una FOTO al informe — la app los lee sola!"),
        ("3️⃣", "Hace la operación",
         "Apunta la hora exacta: cuánto tiempo estuvo el hígado sin sangre"),
        ("4️⃣", "Sigue al paciente después",
         "Cada día apunta cómo va: análisis, si hubo problemas, cuándo se fue a casa"),
        ("5️⃣", "Ve las estadísticas",
         "La app hace gráficas automáticas para aprender y mejorar"),
    ]

    y = 1.5
    colors_paso = [ORANGE, BLUE, GREEN, PURPLE, TEAL]
    for (num, titulo, desc), color in zip(pasos, colors_paso):
        add_rect(s, Inches(0.4), Inches(y), Inches(0.8), Inches(0.9), color)
        add_text_box(s, num, Inches(0.4), Inches(y), Inches(0.8), Inches(0.9),
                     font_size=28, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(s, titulo, Inches(1.4), Inches(y), Inches(11), Inches(0.45),
                     font_size=20, bold=True, color=color)
        add_text_box(s, desc, Inches(1.4), Inches(y+0.42), Inches(11.2), Inches(0.42),
                     font_size=16, color=DARK)
        y += 1.05


def slide_secreto(prs):
    """Slide 6 — El secreto (privacidad)"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xFC, 0xE4, 0xEC))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), RED)
    add_text_box(s, "El gran secreto 🔒",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s, "🕵️",
                 Inches(0.3), Inches(1.5), Inches(2.5), Inches(2.5),
                 font_size=120, align=PP_ALIGN.CENTER)

    add_text_box(s,
        "Los pacientes tienen derecho a que\nsu información sea PRIVADA",
        Inches(3.2), Inches(1.5), Inches(9.5), Inches(1.0),
        font_size=26, bold=True, color=RED)

    secretos = [
        "🙈  La app NUNCA guarda el nombre del paciente",
        "🔢  En su lugar, usa un código secreto: BDT-20260508-001",
        "🗝️   Solo el médico con su PIN puede saber quién es",
        "📵  Los datos NO viajan por internet — están solo en el móvil",
        "🧊  Si alguien roba el móvil, no puede leer nada sin el PIN",
    ]
    y = 2.9
    for item in secretos:
        add_text_box(s, item, Inches(3.2), Inches(y), Inches(9.8), Inches(0.6),
                     font_size=20, color=DARK)
        y += 0.88


def slide_graficas(prs):
    """Slide 7 — Las gráficas"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xE8, 0xEA, 0xF6))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), PURPLE)
    add_text_box(s, "Las gráficas mágicas 📊",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s,
        "¿Recuerdas los dibujos de tu cole que muestran cuántos niños\n"
        "tienen perro, gato o pez? La app hace lo mismo con los trasplantes:",
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
        font_size=22, color=DARK)

    graficas = [
        ("📅", "Cuántos trasplantes\nhay cada mes", BLUE),
        ("🩸", "Cuántos donantes\nson de cada tipo", TEAL),
        ("💊", "Qué enfermedades\ntienen los pacientes", PURPLE),
        ("⚠️", "Qué complicaciones\npasan más", ORANGE),
        ("❤️‍🩹", "Cuántos pacientes\nse recuperaron bien", GREEN),
        ("📈", "Si los resultados\nmejoran cada año", RED),
    ]

    x = 0.4
    for (emoji, texto, color) in graficas:
        add_rect(s, Inches(x), Inches(2.7), Inches(1.9), Inches(2.8), color)
        add_text_box(s, emoji, Inches(x), Inches(2.8), Inches(1.9), Inches(0.8),
                     font_size=44, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(s, texto, Inches(x), Inches(3.6), Inches(1.9), Inches(0.9),
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += 2.1

    add_text_box(s,
        "Con estas gráficas los médicos aprenden qué hacen bien\n"
        "y qué pueden mejorar para salvar más vidas",
        Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
        font_size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


def slide_sin_internet(prs):
    """Slide 8 — Sin internet"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xE0, 0xF7, 0xFA))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), TEAL)
    add_text_box(s, "¡Sin internet funciona igual! ✈️",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s, "📚",
                 Inches(0.3), Inches(1.5), Inches(2.5), Inches(2.5),
                 font_size=120, align=PP_ALIGN.CENTER)

    add_text_box(s,
        "La app es como un libro de papel:\nsiempre funciona, aunque no haya WiFi",
        Inches(3.2), Inches(1.5), Inches(9.5), Inches(1.0),
        font_size=26, bold=True, color=TEAL)

    cosas = [
        ("✅ SÍ puede hacer sin internet",
         ["Registrar un caso nuevo", "Ver estadísticas y gráficas",
          "Buscar pacientes", "Exportar a Excel"],
         GREEN),
        ("⚡ Solo necesita WiFi del hospital para...",
         ["Leer fotos con inteligencia artificial",
          "Enviar el Excel al ordenador del hospital"],
         ORANGE),
    ]

    x = 3.2
    for titulo, items_list, color in cosas:
        add_rect(s, Inches(x), Inches(2.8), Inches(4.6), Inches(3.8), color)
        add_text_box(s, titulo, Inches(x+0.1), Inches(2.85), Inches(4.4), Inches(0.6),
                     font_size=15, bold=True, color=WHITE)
        y2 = 3.55
        for item in items_list:
            add_text_box(s, f"• {item}", Inches(x+0.2), Inches(y2),
                         Inches(4.2), Inches(0.45),
                         font_size=16, color=WHITE)
            y2 += 0.5
        x += 5.0


def slide_resumen(prs):
    """Slide 9 — Resumen"""
    s = blank_slide(prs)
    fill_slide(s, PURPLE)

    add_rect(s, Inches(0), Inches(0), W, Inches(1.8), RGBColor(0x38, 0x00, 0x6B))

    add_text_box(s, "¿Qué hemos aprendido hoy? 🌟",
                 Inches(0.5), Inches(0.3), Inches(12), Inches(1.0),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    resumen = [
        ("🫀", "El hígado es muy importante — a veces hay que cambiarlo"),
        ("📱", "BDTH es la libreta digital del médico de trasplantes"),
        ("🦸 + 🤒", "Guarda datos del donante y del receptor"),
        ("🔒", "Los datos están protegidos con un código secreto"),
        ("📊", "Hace gráficas solas para aprender y mejorar"),
        ("✈️", "Funciona siempre, aunque no haya internet"),
        ("❤️", "Todo esto es para salvar más vidas"),
    ]

    y = 2.0
    for emoji, texto in resumen:
        add_text_box(s, emoji, Inches(0.5), Inches(y), Inches(1.0), Inches(0.7),
                     font_size=28, align=PP_ALIGN.CENTER)
        add_text_box(s, texto, Inches(1.7), Inches(y+0.05), Inches(10.8), Inches(0.6),
                     font_size=21, color=WHITE)
        y += 0.73

    add_text_box(s, "Hospital Gregorio Marañón · Unidad de Trasplante Hepático · Mayo 2026",
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.4),
                 font_size=13, color=PURPLE_L, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def slide_datos_seguros(prs):
    """Slide 10 — Garantia: datos que nunca salen del hospital"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0x1A, 0x23, 0x7E))  # azul oscuro

    add_rect(s, Inches(0), Inches(0), W, Inches(1.5), RGBColor(0x0D, 0x14, 0x4A))
    add_text_box(s, "Garantia de seguridad para el jefe 🛡️",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(1.0),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Diagrama local
    add_rect(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(3.2),
             RGBColor(0x28, 0x36, 0x93))
    add_text_box(s, "RED INTERNA DEL HOSPITAL (WiFi)",
                 Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.5),
                 font_size=14, bold=True, color=RGBColor(0x90, 0xCA, 0xF9),
                 align=PP_ALIGN.CENTER)

    add_text_box(s, "📱\nMovil\ndel medico",
                 Inches(0.8), Inches(2.2), Inches(2.0), Inches(1.5),
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.5), Inches(2.2), Inches(1.2), Inches(0.06),
             RGBColor(0xFF, 0xFF, 0xFF))

    add_text_box(s, "──WiFi──>",
                 Inches(2.9), Inches(2.4), Inches(1.8), Inches(0.5),
                 font_size=18, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

    add_text_box(s, "🖥️\nOrdenador\ndel hospital",
                 Inches(4.8), Inches(2.2), Inches(2.2), Inches(1.5),
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s, "<──resultado──",
                 Inches(7.1), Inches(2.4), Inches(2.0), Inches(0.5),
                 font_size=18, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)

    add_text_box(s, "📱\nMovil\n(formulario\nrelleno)",
                 Inches(9.2), Inches(2.15), Inches(2.2), Inches(1.6),
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s, "✗  NINGUNA CONEXION A INTERNET  ✗",
                 Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
                 font_size=20, bold=True, color=RGBColor(0xFF, 0x80, 0x80),
                 align=PP_ALIGN.CENTER)

    garantias = [
        ("🔒", "NHC nunca en texto plano — convertido en codigo secreto antes de guardarse"),
        ("📵", "Fotos y PDFs procesados en el hospital — jamas viajan a internet"),
        ("🚫", "Sin cuenta en la nube, sin Google, sin Amazon, sin Microsoft Azure"),
        ("🔍", "Codigo fuente disponible para auditoria del Servicio de Informatica"),
    ]
    y = 4.7
    for emoji, texto in garantias:
        add_text_box(s, emoji, Inches(0.5), Inches(y), Inches(0.7), Inches(0.55),
                     font_size=22, align=PP_ALIGN.CENTER)
        add_text_box(s, texto, Inches(1.4), Inches(y+0.05), Inches(11.3), Inches(0.5),
                     font_size=17, color=WHITE)
        y += 0.62


def slide_presupuesto(prs):
    """Slide 11 — Presupuesto"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0xE8, 0xF5, 0xE9))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), GREEN)
    add_text_box(s, "Cuanto cuesta? 💶",
                 Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Tabla visual
    headers = ["Concepto", "Coste"]
    filas = [
        ("Desarrollo de la app", "0 EUR"),
        ("Servidor en la nube", "0 EUR  (no existe)"),
        ("Base de datos", "0 EUR  (en el movil)"),
        ("Licencias de software", "0 EUR  (todo open source)"),
        ("OCR local (Tesseract)", "0 EUR  (gratis y local)"),
        ("Mantenimiento anual", "< 10 horas / ano"),
        ("OCR con IA (opcional)", "< 1 EUR / ano  (50 trasplantes)"),
    ]

    # Header
    add_rect(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.55), GREEN)
    add_text_box(s, "CONCEPTO", Inches(0.5), Inches(1.55), Inches(7.5), Inches(0.4),
                 font_size=15, bold=True, color=WHITE)
    add_text_box(s, "COSTE", Inches(8.2), Inches(1.55), Inches(4.5), Inches(0.4),
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    y = 2.1
    for i, (concepto, coste) in enumerate(filas):
        bg = RGBColor(0xF1, 0xF8, 0xE9) if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.4), Inches(y), Inches(12.5), Inches(0.52), bg)
        add_text_box(s, concepto, Inches(0.5), Inches(y+0.05), Inches(7.5), Inches(0.42),
                     font_size=17, color=DARK)
        color_coste = GREEN if "0 EUR" in coste else (ORANGE if "1 EUR" in coste else DARK)
        add_text_box(s, coste, Inches(8.2), Inches(y+0.05), Inches(4.5), Inches(0.42),
                     font_size=17, bold=("0 EUR" in coste), color=color_coste,
                     align=PP_ALIGN.RIGHT)
        y += 0.52

    # Total
    add_rect(s, Inches(0.4), Inches(y+0.1), Inches(12.5), Inches(0.65),
             RGBColor(0x1B, 0x5E, 0x20))
    add_text_box(s, "TOTAL 5 ANOS", Inches(0.5), Inches(y+0.15), Inches(7.5), Inches(0.5),
                 font_size=20, bold=True, color=WHITE)
    add_text_box(s, "< 5 EUR", Inches(8.2), Inches(y+0.15), Inches(4.5), Inches(0.5),
                 font_size=24, bold=True, color=RGBColor(0x69, 0xF0, 0xAE),
                 align=PP_ALIGN.RIGHT)


def slide_convencer_jefe(prs):
    """Slide 12 — Como convencer al jefe"""
    s = blank_slide(prs)
    fill_slide(s, RGBColor(0x37, 0x47, 0x4F))

    add_rect(s, Inches(0), Inches(0), W, Inches(1.4), RGBColor(0x1C, 0x2C, 0x33))
    add_text_box(s, "Lo que le preocupa al jefe — y la respuesta 🤝",
                 Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9),
                 font_size=33, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    objeciones = [
        (
            "😰  'Los datos sensibles de los PDFs\npueden llegar a la API de Claude'",
            "✅  En modo local (por defecto), los PDFs se procesan\nen el ordenador del hospital. Ningun byte sale.",
            RED, GREEN,
        ),
        (
            "💸  'Cuanto va a costar esto\ncada ano?'",
            "✅  Cero euros de infraestructura.\nSi se usa la IA opcional: menos de 1 euro al ano.",
            ORANGE, GREEN,
        ),
        (
            "🔓  'Como se que solo mi equipo\naccede a los datos?'",
            "✅  Sin backend remoto. Sin cuenta en la nube.\nSolo el medico con su movil y su PIN.",
            RED, GREEN,
        ),
        (
            "📋  'Cumple el RGPD\ny la LOPD?'",
            "✅  Si. Privacidad por diseno (art. 25 RGPD).\nNHC nunca en texto plano. Auditable.",
            ORANGE, GREEN,
        ),
    ]

    y = 1.55
    for objecion, respuesta, c_obj, c_resp in objeciones:
        add_rect(s, Inches(0.3), Inches(y), Inches(6.0), Inches(1.05),
                 RGBColor(0x4A, 0x14, 0x14))
        add_text_box(s, objecion, Inches(0.4), Inches(y+0.05), Inches(5.8), Inches(0.95),
                     font_size=14, color=RGBColor(0xFF, 0xCC, 0xCC))

        add_text_box(s, "→", Inches(6.4), Inches(y+0.3), Inches(0.5), Inches(0.5),
                     font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_rect(s, Inches(7.0), Inches(y), Inches(6.0), Inches(1.05),
                 RGBColor(0x14, 0x3D, 0x1A))
        add_text_box(s, respuesta, Inches(7.1), Inches(y+0.05), Inches(5.8), Inches(0.95),
                     font_size=14, color=RGBColor(0xCC, 0xFF, 0xCC))
        y += 1.2

    add_text_box(s,
        "El codigo fuente esta disponible para auditoria por el Servicio de Informatica del hospital en cualquier momento.",
        Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
        font_size=13, color=RGBColor(0xB0, 0xBE, 0xC5), align=PP_ALIGN.CENTER)


def main():
    prs = new_prs()

    slide_portada(prs)
    slide_que_es_higado(prs)
    slide_que_es_app(prs)
    slide_donante_receptor(prs)
    slide_como_funciona(prs)
    slide_secreto(prs)
    slide_graficas(prs)
    slide_sin_internet(prs)
    slide_resumen(prs)
    slide_datos_seguros(prs)
    slide_presupuesto(prs)
    slide_convencer_jefe(prs)

    out = os.path.join(os.path.dirname(__file__), "BDTH_presentacion_5años.pptx")
    prs.save(out)
    print(f"✅ Guardado en: {out}")


if __name__ == "__main__":
    main()
